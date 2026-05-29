import unittest
from io import BytesIO
from types import SimpleNamespace

from pptx import Presentation
from pypdf import PdfReader

from app.services.formula_export import (
    EXPORT_BLOCK_GAP_PPTX_PT,
    EXPORT_BLOCK_GAP_PT,
    EXPORT_FONT_BODY_PT,
    EXPORT_LINE_GAP_PT,
    EXPORT_LINE_GAP_PPTX_PT,
    compute_formula_layout,
    is_tall_formula,
    layout_rich_text,
    next_baseline_after_block,
    split_text_and_formulas,
    estimate_text_width_pt,
)
from app.services.latex_renderer import render_latex_to_png
from app.services.presentation_export_service import (
    _PPTX_BODY_MAX_LINES,
    _format_numbered_option,
    _paginate_body_lines,
    _pptx_advance_after_line,
    _question_body_paragraphs,
    _simplify_latex_for_export,
    _split_text_and_formulas,
    _visual_line_count,
    _build_pdf,
    _build_pptx,
    _latex_skip_omml_use_plaintext,
)


class TestLatexSimplify(unittest.TestCase):
    def test_cap_symbol_in_simplify_fallback(self):
        self.assertIn("∩", _simplify_latex_for_export(r"A \cap B"))

    def test_split_display_and_inline(self):
        parts = _split_text_and_formulas(r"inline $a$ and $$b=c$$")
        kinds = [(k, d) for k, _c, d in parts]
        self.assertIn(("math", False), kinds)
        self.assertIn(("math", True), kinds)

    def test_numbered_option_prefix(self):
        self.assertEqual(
            _format_numbered_option(1, r"$\frac{1}{2}$"),
            r"1. $\frac{1}{2}$",
        )
        self.assertEqual(
            _format_numbered_option(3, "текст"),
            "3. текст",
        )

    def test_frac_and_dollars(self):
        raw = r"Найдите $\frac{1}{b-a}$ при $x>0$"
        simplified = _simplify_latex_for_export(raw)
        self.assertIn("(1)/(b-a)", simplified)
        self.assertNotIn("\\frac", simplified)
        self.assertNotIn("$", simplified)

    def test_mhchem_skips_omml(self):
        self.assertTrue(_latex_skip_omml_use_plaintext(r"\ce{C6H5Br}"))
        self.assertFalse(_latex_skip_omml_use_plaintext(r"x^2+1"))


class TestPptxPagination(unittest.TestCase):
    def _long_question(self):
        return SimpleNamespace(
            question_text=" ".join(["длинный текст вопроса"] * 40),
            question_type="single_choice",
            answers=[f"вариант {i}" for i in range(12)],
            correct_answers=["вариант 1"],
            explanation=" ".join(["пояснение"] * 30),
        )

    def test_many_lines_split_into_pages(self):
        lines = _question_body_paragraphs(self._long_question(), "teacher")
        pages = _paginate_body_lines(lines)
        self.assertGreater(len(pages), 1)
        for page in pages:
            visual_lines = sum(_visual_line_count(text, level) for text, level in page)
            max_single = max(_visual_line_count(text, level) for text, level in page)
            # одна очень длинная строка может занять целый слайд и превысить лимит
            if max_single <= _PPTX_BODY_MAX_LINES:
                self.assertLessEqual(visual_lines, _PPTX_BODY_MAX_LINES)

    def test_build_pptx_creates_continuation_slides(self):
        quiz = SimpleNamespace(
            title="Тест",
            subject="Математика",
            grade="9",
            difficulty="medium",
        )
        questions = [self._long_question()]
        data, _, _ = _build_pptx(quiz, questions, "teacher")
        prs = Presentation(__import__("io").BytesIO(data))
        titles = [slide.shapes.title.text for slide in prs.slides]
        self.assertIn("Вопрос 1 (продолжение)", titles)

    def test_build_pptx_with_mhchem_opens(self):
        """mhchem (\\ce) не должен давать битый OMML — сlide читается python-pptx."""
        quiz = SimpleNamespace(
            title="Химия",
            subject="Химия",
            grade="10",
            difficulty="hard",
        )
        q = SimpleNamespace(
            question_text="Бромирование бензола?",
            question_type="single_choice",
            answers=[r"$\ce{C6H5Br}$", r"$\ce{C6H6}$", r"$\ce{C6H5CH3}$"],
            correct_answers=[r"$\ce{C6H5Br}$"],
            explanation=r"Продукт $\ce{C6H5Br}$.",
        )
        data, _, _ = _build_pptx(quiz, [q], "teacher")
        prs = Presentation(__import__("io").BytesIO(data))
        self.assertGreaterEqual(len(prs.slides), 2)

    def test_build_pptx_mixed_text_and_formulas_uses_shapes(self):
        """Строки с текстом и формулами рендерятся absolute shapes, не placeholder."""
        quiz = SimpleNamespace(
            title="Множества",
            subject="Алгебра",
            grade="10",
            difficulty="medium",
        )
        q = SimpleNamespace(
            question_text="Какие свойства выполняются для пересечения и объединения множеств?",
            question_type="multiple_choice",
            answers=[
                r"коммутативности ($A \cap B = B \cap A$)",
                r"дистрибутивности ($A \cap (B \cup C) = (A \cap B) \cup (A \cap C)$)",
                "ассоциативности только для пересечения",
            ],
            correct_answers=["вариант 1"],
            explanation=r"См. свойства: $A \cap B = B \cap A$.",
        )
        data, _, _ = _build_pptx(quiz, [q], "teacher")
        prs = Presentation(__import__("io").BytesIO(data))
        body_slide = prs.slides[1]
        shape_count = len(body_slide.shapes)
        # title + body placeholder + absolute textboxes/pictures
        self.assertGreater(shape_count, 8)
        body_tf = body_slide.shapes.placeholders[1].text_frame
        body_text = "\n".join(p.text for p in body_tf.paragraphs if p.text.strip())
        self.assertEqual(body_text, "")


class TestFormulaLayout(unittest.TestCase):
    def test_tall_formula_detection(self):
        self.assertTrue(is_tall_formula(r"\frac{a}{b}"))
        self.assertTrue(is_tall_formula(r"\ce{H2O}"))
        self.assertFalse(is_tall_formula(r"A \cap B"))

    def test_frac_taller_than_inline_symbol(self):
        frac_png = render_latex_to_png(r"\frac{a+b}{a-b}", font_size_pt=11)
        cap_png = render_latex_to_png(r"A \cap B", font_size_pt=11)
        self.assertIsNotNone(frac_png)
        self.assertIsNotNone(cap_png)
        frac_h = compute_formula_layout(
            frac_png, 11, latex=r"\frac{a+b}{a-b}",
        ).height_pt
        cap_h = compute_formula_layout(
            cap_png, 11, latex=r"A \cap B",
        ).height_pt
        self.assertGreater(frac_h, cap_h)

    def test_split_text_and_formulas_matches_presentation(self):
        parts = split_text_and_formulas(r"inline $a$ and $$b=c$$")
        legacy = _split_text_and_formulas(r"inline $a$ and $$b=c$$")
        self.assertEqual(parts, legacy)

    def test_layout_keeps_text_and_formula_on_one_line(self):
        def resolve(latex, display):
            png = render_latex_to_png(latex, font_size_pt=11, display=display)
            if not png:
                return None, None, latex
            return png, compute_formula_layout(png, 11, display=display, latex=latex), None

        layout = layout_rich_text(
            r"Найдите $\frac{2x}{3y}$?",
            max_width_pt=500,
            font_size_pt=11,
            measure_text=lambda s: len(s) * 5.5,
            resolve_formula=resolve,
        )
        self.assertEqual(len(layout.lines), 1)
        kinds = [seg.kind for line in layout.lines for seg in line.segments]
        self.assertIn("text", kinds)
        self.assertIn("math", kinds)

    def test_fraction_options_have_tall_lines(self):
        options = [
            r"- $-\frac{2x}{3y}$",
            r"- $\frac{4x}{6y}$",
            r"- $\frac{2x^2}{3y^2}$",
        ]

        def resolve(latex, display):
            png = render_latex_to_png(latex, font_size_pt=11, display=display)
            if not png:
                return None, None, latex
            return png, compute_formula_layout(png, 11, display=display, latex=latex), None

        heights = []
        for option in options:
            layout = layout_rich_text(
                option,
                max_width_pt=500,
                font_size_pt=11,
                measure_text=lambda s: len(s) * 5.5,
                resolve_formula=resolve,
            )
            self.assertGreaterEqual(len(layout.lines), 1)
            heights.append(layout.lines[0].height_pt)

        self.assertTrue(all(h > EXPORT_FONT_BODY_PT for h in heights))
        self.assertGreater(sum(heights), len(heights) * EXPORT_LINE_GAP_PT)

    def test_pptx_gaps_separate_from_pdf(self):
        self.assertGreater(EXPORT_LINE_GAP_PPTX_PT, EXPORT_LINE_GAP_PT)
        self.assertGreater(EXPORT_BLOCK_GAP_PPTX_PT, EXPORT_BLOCK_GAP_PT)

    def test_pptx_block_gap_exceeds_line_gap(self):
        line_step = _pptx_advance_after_line(
            14.0, has_next=True, next_is_block_start=False,
        )
        block_step = _pptx_advance_after_line(
            14.0, has_next=True, next_is_block_start=True,
        )
        self.assertGreater(block_step, line_step)

    def test_next_baseline_separates_text_options(self):
        options = [
            "текст без формул",
            "ещё один вариант ответа",
            "третий вариант",
        ]
        start_y = 700.0
        y = start_y
        for option in options:
            layout = layout_rich_text(
                option,
                max_width_pt=500,
                font_size_pt=11,
                measure_text=lambda s: len(s) * 5.5,
                resolve_formula=lambda *_: (None, None, None),
            )
            next_y = next_baseline_after_block(y, layout)
            self.assertLess(
                next_y,
                y - layout.lines[0].height_pt,
                "следующий baseline должен учитывать высоту строки и gap",
            )
            y = next_y
        total_drop = start_y - y
        min_expected = sum(
            layout_rich_text(
                opt, 500, 11,
                measure_text=lambda s: len(s) * 5.5,
                resolve_formula=lambda *_: (None, None, None),
            ).lines[0].height_pt + EXPORT_BLOCK_GAP_PT
            for opt in options
        )
        self.assertGreaterEqual(total_drop, min_expected)


class TestPdfLatex(unittest.TestCase):
    def test_render_cap_png(self):
        png = render_latex_to_png(r"A \cap B")
        self.assertIsNotNone(png)
        self.assertGreater(len(png or b""), 200)

    def test_build_pdf_with_fractions(self):
        quiz = SimpleNamespace(
            title="Алгебра",
            subject="Алгебра",
            grade="8",
            difficulty="medium",
        )
        q = SimpleNamespace(
            question_text=r"Найдите $\frac{a+b}{a-b}$ при $a=3$, $b=2$.",
            question_type="single_choice",
            answers=[r"$-\frac{5}{1}$", r"$\frac{1}{5}$", "текст"],
            correct_answers=[r"$\frac{5}{1}$"],
            explanation=r"Подставляем: $\frac{3+2}{3-2}=5$.",
        )
        data, _, media_type = _build_pdf(quiz, [q], "teacher")
        self.assertEqual(media_type, "application/pdf")
        self.assertTrue(data.startswith(b"%PDF"))

    def test_build_pdf_with_inline_latex(self):
        quiz = SimpleNamespace(
            title="Математика",
            subject="Алгебра",
            grade="9",
            difficulty="medium",
        )
        q = SimpleNamespace(
            question_text=r"Чему равно $A \cap B$?",
            question_type="single_choice",
            answers=[r"$A \cap B$", r"$A \cup B$", "текст без формул"],
            correct_answers=[r"$A \cap B$"],
            explanation=r"Пересечение: $A \cap B$.",
        )
        data, _, media_type = _build_pdf(quiz, [q], "teacher")
        self.assertEqual(media_type, "application/pdf")
        self.assertGreater(len(data), 500)
        self.assertTrue(data.startswith(b"%PDF"))

    def test_build_pdf_short_questions_share_page(self):
        quiz = SimpleNamespace(
            title="Компактный тест",
            subject="Математика",
            grade="8",
            difficulty="easy",
        )

        def short_question(number: int):
            return SimpleNamespace(
                question_text=f"Чему равно {number} + 1?",
                question_type="single_choice",
                answers=[str(number), str(number + 1), str(number + 2)],
                correct_answers=[str(number + 1)],
                explanation="Прибавьте единицу.",
            )

        data, _, media_type = _build_pdf(
            quiz,
            [short_question(1), short_question(2), short_question(3)],
            "teacher",
        )
        self.assertEqual(media_type, "application/pdf")
        self.assertEqual(len(PdfReader(BytesIO(data)).pages), 1)


if __name__ == "__main__":
    unittest.main()
