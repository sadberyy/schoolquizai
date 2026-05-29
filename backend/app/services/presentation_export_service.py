import math
import os
import re
import tempfile
from copy import deepcopy
from io import BytesIO
from pathlib import Path
from typing import Literal

from fastapi import HTTPException
from lxml import etree
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.util import Pt

from reportlab.lib.pagesizes import A4
from reportlab.lib.units import mm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas

import matplotlib
matplotlib.use('Agg')

from app.services.formula_export import (
    EXPORT_BLOCK_GAP_PT,
    EXPORT_BLOCK_GAP_PPTX_PT,
    EXPORT_FONT_BODY_PT,
    EXPORT_FONT_LINE_HEIGHT_PT,
    EXPORT_LINE_GAP_PT,
    EXPORT_LINE_GAP_PPTX_PT,
    FormulaLayout,
    LayoutLine,
    compute_formula_layout,
    layout_rich_text,
    next_baseline_after_block,
    pptx_font_size_pt,
    split_text_and_formulas as _split_text_and_formulas,
)
from app.services.latex_renderer import render_latex_to_png, latex_render_batch
from app.db.database import get_db_session
from app.db.models import Question, Quiz

ExportFormat = Literal["pptx", "pdf"]
ExportMode = Literal["teacher", "student"]

_BUNDLED_FONT = Path(__file__).resolve().parent.parent / "assets" / "fonts" / "DejaVuSans.ttf"
_BUNDLED_FONT_BOLD = Path(__file__).resolve().parent.parent / "assets" / "fonts" / "DejaVuSans-Bold.ttf"
_FONT_CANDIDATES = [
    _BUNDLED_FONT,
    Path("C:/Windows/Fonts/arial.ttf"),
    Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
    Path("/usr/share/fonts/TTF/DejaVuSans.ttf"),
    Path("/System/Library/Fonts/Supplemental/Arial.ttf"),
]
_DIFFICULTY_LABELS = {"easy": "Лёгкая", "medium": "Средняя", "hard": "Сложная"}
_TYPE_LABELS = {
    "single_choice": "Один вариант",
    "multiple_choice": "Несколько вариантов",
    "true_false": "Верно / Неверно",
}

_FONT_REGISTERED = False

_PPTX_ROW_HEIGHT_PT = 24.0  # ~14 pt строка + средний gap для оценки пагинации
_PPTX_BODY_WIDTH_CHARS = 60
_PPTX_BODY_MAX_LINES   = 12   # сколько визуальных строк помещается по высоте при 18pt
_PPTX_WRAP_CHARS = 78
_PPTX_FONT_NAME = "Arial"
_PPTX_FONT_REGISTERED = False
_LATEX_MARKER_RE = re.compile(
    r"\\frac|\\sqrt|\$\$?|\\\[|\\\]|\\[a-zA-Z]{2,}"
)
_FRAC_RE = re.compile(r"\\frac\{([^{}]*)\}\{([^{}]*)\}")
_LATEX_SYMBOLS = (
    (r"\\cdot", "·"),
    (r"\\times", "×"),
    (r"\\pm", "±"),
    (r"\\leq", "≤"),
    (r"\\geq", "≥"),
    (r"\\neq", "≠"),
    (r"\\infty", "∞"),
    (r"\\pi", "π"),
    (r"\\alpha", "α"),
    (r"\\beta", "β"),
    (r"\\gamma", "γ"),
    (r"\\delta", "δ"),
    (r"\\theta", "θ"),
    (r"\\lambda", "λ"),
    (r"\\sigma", "σ"),
    (r"\\rightarrow", "→"),
    (r"\\leftarrow", "←"),
    (r"\\cap", "∩"),
    (r"\\cup", "∪"),
    (r"\\subset", "⊂"),
    (r"\\supset", "⊃"),
    (r"\\subseteq", "⊆"),
    (r"\\supseteq", "⊇"),
    (r"\\in", "∈"),
    (r"\\notin", "∉"),
    (r"\\emptyset", "∅"),
    (r"\\forall", "∀"),
    (r"\\exists", "∃"),
)

_A_NS   = "http://schemas.openxmlformats.org/drawingml/2006/main"
_MC_NS  = "http://schemas.openxmlformats.org/markup-compatibility/2006"
_A14_NS = "http://schemas.microsoft.com/office/drawing/2010/main"
_M_NS   = "http://schemas.openxmlformats.org/officeDocument/2006/math"

# mhchem / расширения LaTeX: latex2mathml -> OMML даёт невалидный для PowerPoint XML
# (например \ce и тело формулы разбиваются на отдельные <m:r>).
_MHCHEM_OR_UNSUPPORTED_OMML = re.compile(
    r"\\ce\b|\\cf\b|\\cee\b|\\pu\b|\\bond\b|"
    r"\\x(?:left|right)arrow\b|\\xrightleftharpoons\b|\\xleftrightarrow\b",
    re.IGNORECASE,
)


def _latex_skip_omml_use_plaintext(latex: str) -> bool:
    s = (latex or "").strip()
    return bool(s and _MHCHEM_OR_UNSUPPORTED_OMML.search(s))


def _latex_plaintext_for_slide(latex: str) -> str:
    """Формула для слайда без OMML: убрать обёртки mhchem, остальное упростить."""
    s = (latex or "").strip()
    s = re.sub(r"\\ce\{([^}]*)\}", r"\1", s, flags=re.IGNORECASE)
    s = re.sub(r"\\cf\{([^}]*)\}", r"\1", s, flags=re.IGNORECASE)
    s = re.sub(r"\\cee\{([^}]*)\}", r"\1", s, flags=re.IGNORECASE)
    s = re.sub(r"\\pu\{([^}]*)\}", r"\1", s, flags=re.IGNORECASE)
    return _simplify_latex_for_export(s) or s


def _make_text_run(text: str, font_size_pt: int = 14) -> etree._Element:
    a_r = etree.Element(f"{{{_A_NS}}}r", nsmap={"a": _A_NS})
    rpr = etree.SubElement(a_r, f"{{{_A_NS}}}rPr")
    rpr.set("lang", "ru-RU")
    rpr.set("dirty", "0")
    rpr.set("sz", str(font_size_pt * 100))  # pptx хранит размер в сотых пункта
    a_t = etree.SubElement(a_r, f"{{{_A_NS}}}t")
    a_t.text = text
    return a_r


class _PptxFormulaLayout:
    """Размещает текст и PNG-формулы на слайде с абсолютными координатами."""

    _LEVEL_INDENT_PT = 18
    _TOP_PADDING_PT = 6

    def __init__(self, slide) -> None:
        placeholder = slide.shapes.placeholders[1]
        self.slide = slide
        self.body_left = placeholder.left
        self.body_top = placeholder.top

    def _level_indent(self, level: int):
        return Pt(level * self._LEVEL_INDENT_PT)

    @staticmethod
    def _configure_text_frame(text_frame) -> None:
        text_frame.word_wrap = False
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0
        text_frame.vertical_anchor = MSO_ANCHOR.TOP

    def _add_text_box(
        self,
        left,
        top,
        width_pt: float,
        height_pt: float,
        text: str,
        font_size_pt: float,
    ) -> None:
        if not text:
            return
        box = self.slide.shapes.add_textbox(
            left,
            top,
            Pt(max(width_pt + 4, 12)),
            Pt(max(height_pt, font_size_pt + 2)),
        )
        text_frame = box.text_frame
        self._configure_text_frame(text_frame)
        paragraph = text_frame.paragraphs[0]
        paragraph.text = text
        paragraph.font.size = Pt(font_size_pt)
        paragraph.font.name = _PPTX_FONT_NAME

    def render_absolute_line(
        self,
        line: LayoutLine,
        level: int,
        font_size_pt: float,
        y_offset_pt: float,
    ) -> None:
        indent = self._level_indent(level)
        line_top = self.body_top + Pt(y_offset_pt)
        baseline_top = line.ascent_pt

        index = 0
        while index < len(line.segments):
            seg = line.segments[index]
            if seg.kind in {"text", "fallback"}:
                text_parts: list[str] = []
                x_pt = seg.x_pt
                width_pt = 0.0
                max_ascent = seg.ascent_pt
                while index < len(line.segments):
                    current = line.segments[index]
                    if current.kind not in {"text", "fallback"}:
                        break
                    if current.content:
                        text_parts.append(current.content)
                        width_pt += current.width_pt
                        max_ascent = max(max_ascent, current.ascent_pt)
                    index += 1
                text = "".join(text_parts)
                if text:
                    text_top = line_top + Pt(baseline_top - max_ascent)
                    self._add_text_box(
                        self.body_left + indent + Pt(x_pt),
                        text_top,
                        width_pt,
                        max_ascent + font_size_pt * 0.28,
                        text,
                        font_size_pt,
                    )
                continue

            if seg.kind == "math" and seg.png_bytes and seg.formula_layout:
                fl = seg.formula_layout
                formula_top = line_top + Pt(baseline_top - fl.baseline_offset_pt)
                self.slide.shapes.add_picture(
                    BytesIO(seg.png_bytes),
                    self.body_left + indent + Pt(seg.x_pt),
                    formula_top,
                    height=Pt(fl.height_pt),
                )
            index += 1


def _dejavu_measure_text(text: str, font_size_pt: float) -> float:
    _ensure_pdf_fonts()
    return pdfmetrics.stringWidth(text, "DejaVu", font_size_pt)


def _ensure_pptx_font() -> None:
    global _PPTX_FONT_REGISTERED
    if _PPTX_FONT_REGISTERED:
        return
    pptx_font_candidates = [
        Path("C:/Windows/Fonts/arial.ttf"),
        Path("/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"),
        Path("/System/Library/Fonts/Supplemental/Arial.ttf"),
        *_FONT_CANDIDATES,
    ]
    for candidate in pptx_font_candidates:
        if candidate.is_file() and _is_valid_ttf(candidate):
            try:
                pdfmetrics.registerFont(TTFont(_PPTX_FONT_NAME, str(candidate)))
                _PPTX_FONT_REGISTERED = True
                return
            except Exception:
                continue
    _ensure_pdf_fonts()
    _PPTX_FONT_REGISTERED = True


def _pptx_measure_text(text: str, font_size_pt: float) -> float:
    _ensure_pptx_font()
    if _PPTX_FONT_NAME in pdfmetrics.getRegisteredFontNames():
        return pdfmetrics.stringWidth(text, _PPTX_FONT_NAME, font_size_pt)
    return _dejavu_measure_text(text, font_size_pt)


def _resolve_formula_segment(
    latex: str,
    display: bool,
    font_size_pt: float,
) -> tuple[bytes | None, FormulaLayout | None, str | None]:
    png = render_latex_to_png(latex, font_size_pt=font_size_pt, display=display)
    if png:
        return png, compute_formula_layout(
            png, font_size_pt, display=display, latex=latex,
        ), None
    return None, None, _latex_plaintext_for_slide(latex)


PptxLayoutRow = tuple[LayoutLine, int, float, bool]


def _pptx_advance_after_line(
    line_height_pt: float,
    *,
    has_next: bool,
    next_is_block_start: bool,
) -> float:
    advance = line_height_pt
    if has_next:
        advance += (
            EXPORT_BLOCK_GAP_PPTX_PT
            if next_is_block_start
            else EXPORT_LINE_GAP_PPTX_PT
        )
    return advance


def _layout_paragraph_rows(
    paragraphs: list[tuple[str, int]],
    max_width_pt: float,
    *,
    measure_text,
) -> list[PptxLayoutRow]:
    rows: list[PptxLayoutRow] = []
    for text, level in paragraphs:
        font_size = pptx_font_size_pt(level)
        layout = layout_rich_text(
            text or "",
            max_width_pt - (level * _PptxFormulaLayout._LEVEL_INDENT_PT),
            font_size,
            measure_text=lambda s, fs=font_size: measure_text(s, fs),
            resolve_formula=lambda latex, display, fs=font_size: _resolve_formula_segment(
                latex, display, fs,
            ),
        )
        for line_index, line in enumerate(layout.lines):
            rows.append((line, level, font_size, line_index == 0))
    return rows


def _paginate_layout_rows(
    rows: list[PptxLayoutRow],
    max_height_pt: float | None = None,
) -> list[list[PptxLayoutRow]]:
    if max_height_pt is None:
        max_height_pt = _PPTX_BODY_MAX_LINES * _PPTX_ROW_HEIGHT_PT
    if not rows:
        return [[]]

    pages: list[list[PptxLayoutRow]] = []
    current: list[PptxLayoutRow] = []
    current_height = 0.0

    for index, (line, level, font_size, block_start) in enumerate(rows):
        gap_before = 0.0
        if current:
            gap_before = (
                EXPORT_BLOCK_GAP_PPTX_PT if block_start else EXPORT_LINE_GAP_PPTX_PT
            )
        needed = gap_before + line.height_pt
        if current and current_height + needed > max_height_pt:
            pages.append(current)
            current, current_height = [], 0.0
            needed = line.height_pt
        current.append((line, level, font_size, block_start))
        current_height += needed

    if current:
        pages.append(current)
    return pages


def _prepare_paragraph(paragraph, level: int, font_size_pt: int, line_height_pt: float) -> None:
    p_elem = paragraph._p
    for child in list(p_elem):
        tag = child.tag
        if tag == qn("a:r") or tag == qn("a:br") or tag.endswith("}AlternateContent"):
            p_elem.remove(child)

    p_pr = p_elem.find(qn("a:pPr"))
    if p_pr is None:
        p_pr = etree.SubElement(p_elem, qn("a:pPr"))
        p_elem.insert(0, p_pr)
    p_pr.set("lvl", str(level))

    ln_spc = p_pr.find(qn("a:lnSpc"))
    if ln_spc is None:
        ln_spc = etree.SubElement(p_pr, qn("a:lnSpc"))
    for old in list(ln_spc):
        ln_spc.remove(old)
    spc_pts = etree.SubElement(ln_spc, qn("a:spcPts"))
    spc_pts.set("val", str(int(max(line_height_pt, font_size_pt + 4) * 100)))


def _fill_pptx_body(
    text_frame,
    rows: list[PptxLayoutRow],
    slide,
) -> None:
    text_frame.clear()
    text_frame.word_wrap = False
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    layout = _PptxFormulaLayout(slide)
    y_offset_pt = layout._TOP_PADDING_PT

    for index, (line, level, font_size, _block_start) in enumerate(rows):
        layout.render_absolute_line(line, level, font_size, y_offset_pt)
        has_next = index + 1 < len(rows)
        y_offset_pt += _pptx_advance_after_line(
            line.height_pt,
            has_next=has_next,
            next_is_block_start=rows[index + 1][3] if has_next else False,
        )


def _make_math_run(omml_element: etree._Element, fallback_text: str) -> etree._Element:
    """
    Оборачивает OMML в AlternateContent: PowerPoint покажет формулу,
    старые вьюверы — fallback-текст.
    """
    alt = etree.Element(
        f"{{{_MC_NS}}}AlternateContent",
        nsmap={"mc": _MC_NS, "a14": _A14_NS, "m": _M_NS, "a": _A_NS},
    )
    choice = etree.SubElement(alt, f"{{{_MC_NS}}}Choice", Requires="a14")
    a_r = etree.SubElement(choice, f"{{{_A_NS}}}r")
    rpr = etree.SubElement(a_r, f"{{{_A_NS}}}rPr")
    rpr.set("lang", "ru-RU")
    a14_m = etree.SubElement(a_r, f"{{{_A14_NS}}}m")
    a14_m.append(deepcopy(omml_element))

    fallback = etree.SubElement(alt, f"{{{_MC_NS}}}Fallback")
    fallback.append(_make_text_run(fallback_text))
    return alt


def _contains_latex_markers(text: str | None) -> bool:
    return bool(text and _LATEX_MARKER_RE.search(text))


def _simplify_latex_for_export(text: str | None) -> str:
    if not text:
        return ""
    s = str(text)
    while _FRAC_RE.search(s):
        s = _FRAC_RE.sub(lambda m: f"({m.group(1)})/({m.group(2)})", s)
    s = re.sub(r"\\sqrt\{([^{}]*)\}", r"√\1", s)
    s = re.sub(r"\\text\{([^{}]*)\}", r"\1", s)
    s = re.sub(r"\\left|\\right", "", s)
    for pattern, replacement in _LATEX_SYMBOLS:
        s = re.sub(pattern, replacement, s)
    s = re.sub(r"\$\$?", "", s)
    s = re.sub(r"\\\[|\\\]", "", s)
    s = re.sub(r"\\([a-zA-Z]+)", r" \1 ", s)
    while re.search(r"\{[^{}]*\}", s):
        s = re.sub(r"\{([^{}]*)\}", r"\1", s)
    s = re.sub(r"\s+", " ", s).strip()
    return s


def _wrap_text_lines(text: str, max_chars: int = _PPTX_WRAP_CHARS) -> list[str]:
    text = text.strip()
    if not text:
        return []
    if len(text) <= max_chars:
        return [text]

    words = text.split()
    lines: list[str] = []
    current: list[str] = []
    current_len = 0

    for word in words:
        word_len = len(word)
        extra = word_len if not current else word_len + 1
        if current and current_len + extra > max_chars:
            lines.append(" ".join(current))
            current = [word]
            current_len = word_len
        else:
            current.append(word)
            current_len += extra

    if current:
        lines.append(" ".join(current))
    return lines


def _format_option(option) -> str:
    if isinstance(option, str):
        return _simplify_latex_for_export(option)
    if isinstance(option, dict):
        for key in ("text", "label", "id"):
            val = option.get(key)
            if val is not None and val != "":
                return _simplify_latex_for_export(str(val))
        return ""
    return _simplify_latex_for_export(str(option))


def _format_option_raw(option) -> str:
    if isinstance(option, str):
        return option
    if isinstance(option, dict):
        return str(
            (option or {}).get("text") or (option or {}).get("label") or (option or {}).get("id") or ""
        )
    return str(option)


def _format_numbered_option(index: int, option) -> str:
    return f"{index}. {_format_option_raw(option)}"


def _format_answers(correct_answers) -> str:
    if correct_answers is None:
        return ""
    if isinstance(correct_answers, list):
        return ", ".join(str(item) for item in correct_answers)
    return str(correct_answers)


def _is_valid_ttf(path: Path) -> bool:
    try:
        signature = path.read_bytes()[:4]
    except OSError:
        return False
    return signature in {b"\x00\x01\x00\x00", b"OTTO", b"ttcf"}


def _resolve_pdf_font_path() -> Path:
    for candidate in _FONT_CANDIDATES:
        if candidate.is_file() and _is_valid_ttf(candidate):
            return candidate
    raise HTTPException(
        status_code=500,
        detail=(
            "Не найден шрифт для PDF. Положите DejaVuSans.ttf в app/assets/fonts "
            "или установите системный шрифт (Arial / DejaVu)."
        ),
    )


def _ensure_pdf_fonts() -> None:
    global _FONT_REGISTERED
    if _FONT_REGISTERED:
        return

    font_path = _resolve_pdf_font_path()
    pdfmetrics.registerFont(TTFont("DejaVu", str(font_path)))

    bold_path = _BUNDLED_FONT_BOLD if _BUNDLED_FONT_BOLD.is_file() and _is_valid_ttf(_BUNDLED_FONT_BOLD) else font_path
    pdfmetrics.registerFont(TTFont("DejaVu-Bold", str(bold_path)))
    _FONT_REGISTERED = True


def _safe_filename(title: str, extension: str) -> str:
    """Создаёт имя файла"""
    base = re.sub(r'[<>:"/\\|?*]', "", title).strip() or "quiz"
    base = base[:80]
    return f"{base}.{extension}"


# Загрузка викторины
def load_quiz_with_questions(quiz_id: str) -> tuple[Quiz, list[Question]]:
    with get_db_session() as session:
        quiz = session.query(Quiz).filter(Quiz.id == quiz_id).first()
        if not quiz:
            raise HTTPException(status_code=404, detail="Викторина не найдена")

        questions = (
            session.query(Question)
            .filter(Question.quiz_id == quiz_id)
            .order_by(Question.order_idx)
            .all()
        )
        if not questions:
            raise HTTPException(status_code=400, detail="Викторина не содержит вопросов")
        session.expunge(quiz)
        for question in questions:
            session.expunge(question)

        return quiz, questions


def export_quiz(
    quiz_id: str,
    export_format: ExportFormat = "pptx",
    mode: ExportMode = "teacher",
) -> tuple[bytes, str, str]:
    quiz, questions = load_quiz_with_questions(quiz_id)
    with latex_render_batch():
        if export_format == "pptx":
            return _build_pptx(quiz, questions, mode)
        if export_format == "pdf":
            return _build_pdf(quiz, questions, mode)
    raise HTTPException(status_code=400, detail="Формат должен быть pptx или pdf")


def _quiz_subtitle(quiz: Quiz) -> str:
    parts = []
    if quiz.subject:
        parts.append(f"Предмет: {quiz.subject}")
    if quiz.grade:
        parts.append(f"Класс: {quiz.grade}")
    if quiz.difficulty:
        label = _DIFFICULTY_LABELS.get(quiz.difficulty, quiz.difficulty)
        parts.append(f"Сложность: {label}")
    return " · ".join(parts)


def _visual_line_count(text: str, level: int) -> int:
    # с поправкой на отступ для уровня
    width = _PPTX_BODY_WIDTH_CHARS - (4 if level > 0 else 0)
    if not text:
        return 1
    return max(1, math.ceil(len(text) / width))


def _paginate_expanded_lines(
    rows: list[PptxLayoutRow],
    max_lines: int = _PPTX_BODY_MAX_LINES,
) -> list[list[PptxLayoutRow]]:
    return _paginate_layout_rows(
        rows,
        max_height_pt=max_lines * _PPTX_ROW_HEIGHT_PT,
    )


def _paginate_body_lines(
    lines: list[tuple[str, int]],
    max_lines: int = _PPTX_BODY_MAX_LINES,
) -> list[list[tuple[str, int]]]:
    if not lines:
        return [[]]

    pages: list[list[tuple[str, int]]] = []
    current: list[tuple[str, int]] = []
    current_visual = 0

    for text, level in lines:
        needed = _visual_line_count(text, level)
        # если одна «строка» сама по себе больше страницы — кладём её отдельно
        if needed >= max_lines and current:
            pages.append(current)
            current, current_visual = [], 0
        if current_visual + needed > max_lines and current:
            pages.append(current)
            current, current_visual = [], 0
        current.append((text, level))
        current_visual += needed

    if current:
        pages.append(current)
    return pages


def _question_body_paragraphs(question: Question, mode: ExportMode) -> list[tuple[str, int]]:
    paragraphs: list[tuple[str, int]] = []
    paragraphs.append((question.question_text or "", 0))

    type_label = _TYPE_LABELS.get(question.question_type, question.question_type)
    paragraphs.append((f"Тип: {type_label}", 0))

    for index, option in enumerate(question.answers or [], start=1):
        paragraphs.append((_format_numbered_option(index, option), 1))

    if mode == "teacher":
        answer = _format_answers(question.correct_answers)
        paragraphs.append((f"Правильный ответ: {answer}", 0))
        if question.explanation:
            paragraphs.append((f"Пояснение: {question.explanation}", 0))
    return paragraphs


def _build_pptx(
    quiz: Quiz,
    questions: list[Question],
    mode: ExportMode,
) -> tuple[bytes, str, str]:
    prs = Presentation()

    # Титульный слайд
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = quiz.title or "Викторина"

    subtitle_parts = [_quiz_subtitle(quiz)]
    subtitle = " · ".join(part for part in subtitle_parts if part)
    if subtitle and len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = subtitle

    max_width_pt = float(prs.slide_width) / 12700 * 0.82

    # Слайды вопросов
    for index, question in enumerate(questions, start=1):
        paragraphs = _question_body_paragraphs(question, mode)
        rows = _layout_paragraph_rows(
            paragraphs,
            max_width_pt,
            measure_text=_pptx_measure_text,
        )
        pages = _paginate_expanded_lines(rows)
        for part_idx, page in enumerate(pages):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = (
                f"Вопрос {index}" if part_idx == 0 else f"Вопрос {index} (продолжение)"
            )
            _fill_pptx_body(slide.shapes.placeholders[1].text_frame, page, slide)

    buffer = BytesIO()
    prs.save(buffer)
    filename = _safe_filename(quiz.title, "pptx")
    media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    return buffer.getvalue(), filename, media_type


# PDF экспорт
def _clean_latex_text(text: str) -> str:
    """
    Очищает текст от артефактов:
    - Заменяем двойные скобки [[]] на []
    - Заменяет \\ на \
    - Убирает [ ] вокруг LaTeX-формул
    """
    text = text.replace('\\\\', '\\')
    text = text.replace('[[', '[').replace(']]', ']')
    text = re.sub(r'\[\s*(\$.+?\$)\s*\]', r'\1', text)

    return text


def _draw_rich_text_block(
    c: canvas.Canvas,
    text: str,
    x: float,
    y: float,
    max_width: float,
    font_name: str = "DejaVu",
    font_size: float = EXPORT_FONT_BODY_PT,
) -> float:
    """Рисует блок текста+формул; возвращает baseline для следующего блока."""
    cleaned = _clean_latex_text(text or "")

    def measure(token: str) -> float:
        c.setFont(font_name, font_size)
        return c.stringWidth(token, font_name, font_size)

    layout = layout_rich_text(
        cleaned,
        max_width,
        font_size,
        measure_text=measure,
        resolve_formula=lambda latex, display: _resolve_formula_segment(
            latex, display, font_size,
        ),
    )
    if not layout.lines:
        return y - EXPORT_FONT_LINE_HEIGHT_PT - EXPORT_BLOCK_GAP_PT

    baseline_y = y
    for line_index, line in enumerate(layout.lines):
        for seg in line.segments:
            sx = x + seg.x_pt
            if seg.kind in {"text", "fallback"}:
                c.setFont(font_name, font_size)
                c.drawString(sx, baseline_y, seg.content)
                continue

            if not seg.png_bytes or not seg.formula_layout:
                continue

            fl = seg.formula_layout
            tmp_path = None
            try:
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                    tmp.write(seg.png_bytes)
                    tmp_path = tmp.name
                c.drawImage(
                    tmp_path,
                    sx,
                    baseline_y - fl.height_pt + fl.baseline_offset_pt,
                    width=fl.width_pt,
                    height=fl.height_pt,
                )
            finally:
                if tmp_path and os.path.exists(tmp_path):
                    try:
                        os.unlink(tmp_path)
                    except OSError:
                        pass

        if line_index + 1 < len(layout.lines):
            baseline_y -= line.height_pt + EXPORT_LINE_GAP_PT

    return next_baseline_after_block(y, layout)


_PDF_QUESTION_GAP_PT = 24.0
_PDF_MIN_QUESTION_START_Y_PT = 60.0


def _pdf_break_if_low(
    c: canvas.Canvas,
    y: float,
    height: float,
    margin: float,
    *,
    min_y: float = _PDF_MIN_QUESTION_START_Y_PT,
) -> float:
    if y < margin + min_y:
        c.showPage()
        return height - margin
    return y


def _build_pdf(quiz: Quiz, questions: list[Question], mode: ExportMode) -> tuple[bytes, str, str]:
    _ensure_pdf_fonts()

    buf = BytesIO()
    c = canvas.Canvas(buf, pagesize=A4)
    width, height = A4

    margin = 20 * mm
    max_text_width = width - 2 * margin

    x = margin
    y = height - margin

    # титульник
    c.setFont("DejaVu-Bold", 18)
    c.drawString(x, y, quiz.title or "Викторина")
    y -= 25

    subtitle = _quiz_subtitle(quiz)
    if subtitle:
        c.setFont("DejaVu", 12)
        c.drawString(x, y, subtitle)
        y -= 30

    for index, question in enumerate(questions, start=1):
        if index > 1:
            y -= _PDF_QUESTION_GAP_PT
        y = _pdf_break_if_low(c, y, height, margin)

        # номер вопроса
        c.setFont("DejaVu-Bold", 14)
        c.drawString(x, y, f"Вопрос {index}")
        y -= 22

        # текст вопроса
        y = _draw_rich_text_block(
            c, question.question_text or "", x, y, max_text_width,
        )
        y -= 8

        # тип вопроса
        type_label = _TYPE_LABELS.get(question.question_type, question.question_type)
        c.setFont("DejaVu", 10)
        c.drawString(x, y, f"Тип: {type_label}")
        y -= 18

        # варианты ответов
        for index, option in enumerate(question.answers or [], start=1):
            if y < 30:
                c.showPage()
                y = height - margin
            y = _draw_rich_text_block(
                c, _format_numbered_option(index, option), x, y, max_text_width,
            )

        # правильные ответы и пояснение
        if mode == "teacher":
            y -= 8
            if y < 40:
                c.showPage()
                y = height - margin

            c.setFont("DejaVu-Bold", EXPORT_FONT_BODY_PT)
            c.drawString(x, y, "Правильный ответ:")
            y -= 16

            if isinstance(question.correct_answers, list):
                for answer in question.correct_answers:
                    if y < 20:
                        c.showPage()
                        y = height - margin
                    y = _draw_rich_text_block(c, str(answer), x, y, max_text_width)
            else:
                y = _draw_rich_text_block(
                    c, str(question.correct_answers or ""), x, y, max_text_width,
                )

            if question.explanation:
                y -= 8
                if y < 40:
                    c.showPage()
                    y = height - margin
                c.setFont("DejaVu-Bold", EXPORT_FONT_BODY_PT)
                c.drawString(x, y, "Пояснение:")
                y -= 16
                y = _draw_rich_text_block(
                    c, question.explanation or "", x, y, max_text_width,
                )

    c.save()
    buf.seek(0)
    return buf.getvalue(), _safe_filename(quiz.title, "pdf"), "application/pdf"
