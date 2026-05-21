import os
import uuid
import re
from dataclasses import dataclass, field
from typing import Any, Optional, Dict

import fitz
import pdfplumber
import pytesseract
from PIL import Image
from io import BytesIO
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from sqlalchemy import create_engine, Column, String, Text, Integer, DateTime, ForeignKey, JSON
from sqlalchemy.orm import declarative_base, sessionmaker, relationship
from datetime import datetime


# DB MODELS

Base = declarative_base()


class Document(Base):
    __tablename__ = "documents"

    id = Column(String, primary_key=True, default=lambda: str(uuid.uuid4()))
    filename = Column(String, nullable=False)
    file_type = Column(String, nullable=False)
    stored_path = Column(String, nullable=False)
    created_at = Column(DateTime, default=datetime.utcnow)
    meta = Column(JSON, default=dict)

    blocks = relationship("Block", back_populates="document")


class Block(Base):
    __tablename__ = "blocks"

    id = Column(String, primary_key=True, default=lambda: str(uuid.uuid4()))
    document_id = Column(String, ForeignKey("documents.id"), nullable=False)
    block_type = Column(String, nullable=False)
    content = Column(Text, nullable=False)

    page_num = Column(Integer, nullable=True)
    slide_num = Column(Integer, nullable=True)
    order_idx = Column(Integer, default=0)
    image_path = Column(String, nullable=True)
    meta = Column(JSON, default=dict)

    document = relationship("Document", back_populates="blocks")


# DATA STRUCTURE

@dataclass
class ExtractedBlock:
    block_type: str
    content: str
    page_num: Optional[int] = None
    slide_num: Optional[int] = None
    order_idx: int = 0
    image_path: Optional[str] = None
    bbox: Optional[Dict[str, float]] = None          # <-- координаты области
    meta: dict[str, Any] = field(default_factory=dict)


# UTILS

def clean_text(text: str) -> str:
    if not text:
        return ""
    text = text.replace("\x00", " ")
    text = re.sub(r"-\n(\w)", r"\1", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]+\n", "\n", text)
    return text.strip()


def save_original_file(src_path: str, storage_dir: str = "storage") -> str:
    os.makedirs(storage_dir, exist_ok=True)
    ext = os.path.splitext(src_path)[1].lower()
    new_name = f"{uuid.uuid4()}{ext}"
    dst_path = os.path.join(storage_dir, new_name)
    with open(src_path, "rb") as fsrc, open(dst_path, "wb") as fdst:
        fdst.write(fsrc.read())
    return dst_path


def save_block_image(img: Image.Image, storage_dir: str = "storage/images") -> str:
    os.makedirs(storage_dir, exist_ok=True)
    filename = f"{uuid.uuid4()}.png"
    path = os.path.join(storage_dir, filename)
    img.save(path, format="PNG")
    return path


def ocr_pil_image(img: Image.Image, lang: str = "rus+eng+equ") -> str:
    """
    Распознаёт текст с изображения с помощью Tesseract.
    По умолчанию использует русский + английский.
    Для формул можно передавать lang="rus+eng+equ".
    """
    if img is None:
        return ""

    try:
        text = pytesseract.image_to_string(img, lang=lang)
        return clean_text(text)
    except Exception as e:
        print(f"[OCR ERROR] {e}")
        return ""


def looks_like_formula(text: str) -> bool:
    formula_chars = set("∫∑∏√∂∆∞≠≤≥≈≡±×÷²³⁴αβγδεθλμπσφω")
    math_symbols = sum(1 for c in text if c in formula_chars)
    return math_symbols > 2 or bool(re.search(r"[a-zA-Z]\s*[=^_]\s*[\d\w]", text))


def looks_like_chemistry(text: str) -> bool:
    chem_patterns = [r"[A-Z][a-z]?\d*", r"→|⇌|↔", r"[₀-₉⁰-⁹]"]
    return any(re.search(p, text) for p in chem_patterns) and len(text) < 300


# EXTRACTORS

def extract_txt(path: str) -> list[ExtractedBlock]:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        text = clean_text(f.read())
    return [ExtractedBlock(block_type="text", content=text)]


def extract_image(path: str) -> list[ExtractedBlock]:
    img = Image.open(path)
    text = ocr_pil_image(img)
    block_type = "image_ocr" if text else "image"
    return [ExtractedBlock(block_type=block_type, content=text or "[IMAGE_NO_TEXT]", image_path=path)]


def table_to_markdown(table: list[list[str]]) -> str:
    """Конвертирует список списков в Markdown-таблицу."""
    md_rows = []
    if not table:
        return ""

    # Очистка ячеек и замена None на пустые строки
    clean_table = [[str(cell or "").strip() for cell in row] for row in table]

    # Header
    md_rows.append("| " + " | ".join(clean_table[0]) + " |")
    # Separator
    md_rows.append("|" + "---|" * len(clean_table[0]))
    # Body
    for row in clean_table[1:]:
        md_rows.append("| " + " | ".join(row) + " |")
    return "\n".join(md_rows)


def extract_pdf(path: str) -> list[ExtractedBlock]:
    all_blocks: list[ExtractedBlock] = []

    try:
        doc = fitz.open(path)
    except Exception as e:
        return [ExtractedBlock(block_type="note", content=f"[PDF_LOAD_ERROR] {str(e)}")]

    with pdfplumber.open(path) as pdf:
        for page_idx, (page_fitz, page_plumber) in enumerate(zip(doc, pdf.pages), start=1):
            page_elements: list[ExtractedBlock] = []

            # 1. Текстовые блоки (параграфы)
            text_blocks = page_fitz.get_text("blocks")
            for b in text_blocks:
                text = clean_text(b[4])
                if not text or len(text) < 8:
                    continue

                bbox = {"x0": b[0], "y0": b[1], "x1": b[2], "y1": b[3]}
                block_type = "formula" if looks_like_formula(text) else \
                             "chemistry" if looks_like_chemistry(text) else "text"

                page_elements.append(ExtractedBlock(
                    block_type=block_type, content=text, page_num=page_idx, bbox=bbox
                ))

            # 2. Таблицы (через pdfplumber)
            for table_data in page_plumber.extract_tables() or []:
                if not table_data:
                    continue

                table_markdown = table_to_markdown(table_data)

                # Простой bbox таблицы (вся страница — fallback)
                bbox = {
                    "x0": 0, "y0": 0,
                    "x1": page_plumber.width,
                    "y1": page_plumber.height
                }

                page_elements.append(ExtractedBlock(
                    block_type="table",
                    content=table_markdown,
                    page_num=page_idx,
                    bbox=bbox
                ))

            # 3. Изображения
            for img in page_fitz.get_images(full=True):
                xref = img[0]
                try:
                    base_image = doc.extract_image(xref)
                    pil_image = Image.open(BytesIO(base_image["image"]))
                    img_path = save_block_image(pil_image)
                    ocr_text = ocr_pil_image(pil_image)

                    block_type = "formula" if looks_like_formula(ocr_text) else \
                                 "image_ocr" if ocr_text else "image"
                    content = ocr_text if ocr_text else "[IMAGE]"

                    # Простой bbox (всё изображение на странице)
                    rect = page_fitz.rect
                    bbox = {"x0": rect.x0, "y0": rect.y0, "x1": rect.x1, "y1": rect.y1}

                    page_elements.append(ExtractedBlock(
                        block_type=block_type,
                        content=content,
                        page_num=page_idx,
                        image_path=img_path,
                        bbox=bbox
                    ))
                except Exception:
                    continue

            # 4. Сортировка по вертикали
            page_elements.sort(key=lambda el: el.bbox["y0"] if el.bbox else 9999)
            all_blocks.extend(page_elements)

    doc.close()

    # Финальная нумерация
    for i, block in enumerate(all_blocks):
        block.order_idx = i

    return all_blocks


def extract_pptx(path: str) -> list[ExtractedBlock]:
    blocks: list[ExtractedBlock] = []

    try:
        prs = Presentation(path)
    except Exception as e:
        return [ExtractedBlock(block_type="note", content=f"[PPTX_LOAD_ERROR] {str(e)}")]

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes):
            try:
                if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                    text = clean_text(shape.text_frame.text)
                    if text:
                        block_type = "formula" if looks_like_formula(text) else \
                                     "chemistry" if looks_like_chemistry(text) else "text"
                        blocks.append(ExtractedBlock(
                            block_type=block_type,
                            content=text,
                            slide_num=slide_idx,
                            meta={"shape_idx": shape_idx, "source": "pptx_text"}
                        ))

                elif getattr(shape, "has_table", False) and shape.has_table:
                    table = shape.table
                    rows = [clean_text(table.cell(r, c).text) for r in range(len(table.rows))
                            for c in range(len(table.columns))]
                    blocks.append(ExtractedBlock(
                        block_type="table",
                        content=clean_text("\n".join(rows)),
                        slide_num=slide_idx,
                        meta={"shape_idx": shape_idx, "source": "pptx_table"}
                    ))

                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img = Image.open(BytesIO(shape.image.blob))
                    img_path = save_block_image(img)
                    text = ocr_pil_image(img)

                    block_type = "formula" if looks_like_formula(text) or looks_like_chemistry(text) else "image"

                    blocks.append(ExtractedBlock(
                        block_type=block_type,
                        content=text or "[IMAGE]",
                        slide_num=slide_idx,
                        image_path=img_path,
                        meta={"shape_idx": shape_idx, "source": "pptx_image"}
                    ))

            except Exception as e:
                blocks.append(ExtractedBlock(
                    block_type="note",
                    content=f"[SHAPE_ERROR] {str(e)}",
                    slide_num=slide_idx
                ))

    return blocks


# ROUTER + INGEST

def extract_file(path: str) -> list[ExtractedBlock]:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".txt":
        return extract_txt(path)
    elif ext == ".pdf":
        return extract_pdf(path)
    elif ext == ".pptx":
        return extract_pptx(path)
    elif ext in [".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff"]:
        return extract_image(path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def ingest_file(file_path: str, session, storage_dir: str = "storage") -> Document:
    stored_path = None

    try:
        stored_path = save_original_file(file_path, storage_dir)
        filename = os.path.basename(file_path)
        ext = os.path.splitext(file_path)[1].lower()

        doc = Document(
            filename=filename,
            file_type=ext,
            stored_path=stored_path,
            meta={"original_name": filename}
        )
        session.add(doc)
        session.flush()

        try:
            blocks = extract_file(file_path)
        except Exception as e:
            blocks = [ExtractedBlock(block_type="note", content=f"[EXTRACTION_ERROR] {str(e)}")]

        for i, block in enumerate(blocks):
            meta = block.meta or {}
            if block.bbox:
                meta["bbox"] = block.bbox

            db_block = Block(
                document_id=doc.id,
                block_type=block.block_type,
                content=clean_text(block.content),
                page_num=block.page_num,
                slide_num=block.slide_num,
                order_idx=i,
                image_path=block.image_path,
                meta=meta
            )
            session.add(db_block)

        session.commit()
        return doc

    except Exception as e:
        session.rollback()
        if stored_path and os.path.exists(stored_path):
            try:
                os.remove(stored_path)
            except:
                pass
        raise RuntimeError(f"Ошибка при загрузке файла {file_path}: {str(e)}") from e


def init_db(db_url: str = "sqlite:///quiz_ingestion.db"):
    engine = create_engine(db_url)
    Base.metadata.create_all(engine)
    return sessionmaker(bind=engine)