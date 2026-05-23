from io import BytesIO
from typing import List, Tuple

from docx import Document
from pypdf import PdfReader
from pptx import Presentation

from app.schemas.material import SourceFragment


class MaterialService:
    def extract_text_from_txt(self, content: bytes) -> str:
        return content.decode("utf-8", errors="ignore").strip()

    def extract_text_from_pdf(self, content: bytes) -> List[SourceFragment]:
        reader = PdfReader(BytesIO(content))
        fragments = []

        for page_idx, page in enumerate(reader.pages, start=1):
            text = page.extract_text()
            if not text or not text.strip():
                continue

            raw_parts = [part.strip() for part in text.split("\n") if part.strip()]
            chunks = []
            current_chunk = ""

            for part in raw_parts:
                if len(current_chunk) + len(part) + 1 <= 600:
                    current_chunk = f"{current_chunk}\n{part}".strip()
                else:
                    if current_chunk:
                        chunks.append(current_chunk)
                    current_chunk = part

            if current_chunk:
                chunks.append(current_chunk)

            for chunk_idx, chunk_text in enumerate(chunks, start=1):
                fragments.append(
                    SourceFragment(
                        fragment_id=f"pdf_page_{page_idx}_chunk_{chunk_idx}",
                        source_type="pdf",
                        source_name=f"page_{page_idx}_chunk_{chunk_idx}",
                        text=chunk_text
                    )
                )

        return fragments

    def extract_text_from_pptx(self, content: bytes) -> List[SourceFragment]:
        prs = Presentation(BytesIO(content))
        fragments = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            parts = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    parts.append(shape.text.strip())

            slide_text = "\n".join(parts).strip()
            if slide_text:
                fragments.append(
                    SourceFragment(
                        fragment_id=f"pptx_slide_{slide_idx}",
                        source_type="pptx",
                        source_name=f"slide_{slide_idx}",
                        text=slide_text
                    )
                )

        return fragments

    def extract_text_from_docx(self, content: bytes) -> List[SourceFragment]:
        document = Document(BytesIO(content))
        fragments = []

        parts = []

        for paragraph in document.paragraphs:
            if paragraph.text and paragraph.text.strip():
                parts.append(paragraph.text.strip())

        docx_text = "\n".join(parts).strip()

        if docx_text:
            fragments.append(
                SourceFragment(
                    fragment_id="docx_1",
                    source_type="docx",
                    source_name="document",
                    text=docx_text
                )
            )

        return fragments

    def extract_fragments(self, filename: str, content: bytes) -> Tuple[str, List[SourceFragment]]:
        lower_name = filename.lower()

        if lower_name.endswith(".txt"):
            text = self.extract_text_from_txt(content)
            fragments = []

            if text:
                fragments.append(
                    SourceFragment(
                        fragment_id="txt_1",
                        source_type="txt",
                        source_name=filename,
                        text=text
                    )
                )

            return "txt", fragments

        if lower_name.endswith(".pdf"):
            return "pdf", self.extract_text_from_pdf(content)

        if lower_name.endswith(".pptx"):
            return "pptx", self.extract_text_from_pptx(content)

        if lower_name.endswith(".docx"):
            return "docx", self.extract_text_from_docx(content)

        if lower_name.endswith((".png", ".jpg", ".jpeg", ".webp")):
            return "image", []

        raise ValueError(
            "Unsupported file format. Allowed: .txt, .pdf, .pptx, .docx, .png, .jpg, .jpeg, .webp"
        )

    def has_too_little_text(self, fragments: List[SourceFragment]) -> bool:
        total_length = sum(len(fragment.text.strip()) for fragment in fragments if fragment.text)
        non_empty_count = sum(1 for fragment in fragments if fragment.text and fragment.text.strip())
        return total_length < 300 or non_empty_count < 1

    def merge_fragments(self, manual_text: str | None, file_fragments: List[SourceFragment]) -> List[SourceFragment]:
        fragments = []

        cleaned_manual = (manual_text or "").strip()
        if cleaned_manual.lower() in {"string", "source_text", "null", "none"}:
            cleaned_manual = ""

        if cleaned_manual:
            fragments.append(
                SourceFragment(
                    fragment_id="manual_1",
                    source_type="manual_text",
                    source_name="teacher_input",
                    text=cleaned_manual
                )
            )

        fragments.extend(file_fragments)
        return fragments

    def build_combined_context(self, fragments: List[SourceFragment], max_chars: int = 6000) -> str:
        if not fragments:
            return ""

        parts = []

        for fragment in fragments:
            parts.append(
                f"[fragment_id={fragment.fragment_id}; source_type={fragment.source_type}; source_name={fragment.source_name}]\n{fragment.text}"
            )

        combined = "\n\n".join(parts).strip()
        return combined[:max_chars]


material_service = MaterialService()